from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "submission" / "presentation_deck"
ART = ROOT / "docs" / "demo" / "artifacts"

PARCHMENT = RGBColor(0xF5, 0xF4, 0xED)
IVORY = RGBColor(0xFA, 0xF9, 0xF5)
BRAND = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK = RGBColor(0x14, 0x14, 0x13)
DARK_WARM = RGBColor(0x3D, 0x3D, 0x3A)
OLIVE = RGBColor(0x50, 0x4E, 0x49)
STONE = RGBColor(0x6B, 0x6A, 0x64)
BORDER = RGBColor(0xE8, 0xE6, 0xDC)
RED = RGBColor(0xB8, 0x22, 0x17)
GREEN = RGBColor(0x1D, 0x6B, 0x57)

SERIF = "Georgia"
MONO = "Menlo"


SLIDES = [
    {
        "kind": "cover",
        "title": "ClearPath Recovery",
        "subtitle": "No service case closes until authoritative evidence proves the customer is actually restored",
        "meta": "UiPath AgentHack 2026 · Maestro Case · Governed Telecom Recovery",
        "notes": "Open with the business risk: systems can look resolved while the customer still lacks service.",
    },
    {
        "kind": "problem",
        "eyebrow": "01 · BUSINESS PROBLEM",
        "title": "The riskiest case is the one that looks resolved everywhere except the customer’s service.",
        "lead": "CRM, orders, billing, and support notes can all look green while authoritative service evidence is missing, stale, or contradictory.",
        "cards": [
            ("Wrongful closure", "The customer still has no service, but the case is marked resolved."),
            ("Repeat contact", "The next call reopens the same work with worse SLA pressure."),
            ("Audit gap", "Reviewers cannot reconstruct why the case was closed."),
        ],
        "callout": "The project is not a generic governance platform. It is a concrete telecom exception workflow.",
    },
    {
        "kind": "architecture",
        "eyebrow": "02 · ARCHITECTURE THESIS",
        "title": "Agents interpret. Policy decides. Maestro routes. Humans own exceptions.",
        "steps": [
            ("01", "Evidence", "CRM, billing, telemetry, inventory, dispatch, notes"),
            ("02", "Agent", "Turns ambiguous text into structured signals"),
            ("03", "Policy", "Blocks unsafe closure and chooses the allowed route"),
            ("04", "Maestro Case", "Carries stages, tasks, incidents, and lifecycle"),
            ("05", "Human review", "Owns high-impact contradictions and remediation"),
        ],
        "callout": "The LLM never closes cases, overrides policy, or mutates production rules.",
    },
    {
        "kind": "casefit",
        "eyebrow": "03 · WHY MAESTRO CASE",
        "title": "This is Case work because the route emerges as evidence changes.",
        "lead": "The workflow is not a fixed BPMN-style happy path. It is long-running exception handling with changing evidence, SLA pressure, human accountability, and audit requirements.",
        "cards": [
            ("Unpredictable path", "Missing telemetry, stale signals, contradiction, and human review are decided after evidence arrives."),
            ("Human-owned risk", "High-impact contradictions become reviewer tasks with structured return, not silent automation."),
            ("Durable audit", "The case must preserve raw agent recommendation, final policy route, reviewer action, and policy versions."),
        ],
        "callout": "ClearPath Recovery fits Track 1: dynamic exception-heavy casework with agents, automation, and humans in charge at key decision points.",
    },
    {
        "kind": "map",
        "eyebrow": "04 · UIPATH PLATFORM MAP",
        "title": "UiPath is the orchestration and governance boundary.",
        "rows": [
            ("Maestro Case", "Dynamic exception-heavy case lifecycle, stage routing, incidents, recovery path."),
            ("Action Center", "Human task lifecycle, assignment, completion, reviewer action/comment."),
            ("Data Fabric / Orchestrator", "Durable audit bundle, package/process/version readback, full payload proof."),
            ("Test Manager", "Manual execution evidence for E-001 through E-009 eval coverage."),
            ("UiPath CLI", "Repeatable package, task, Data Fabric, Test Manager, and validation readback."),
        ],
        "callout": "Custom evidence packets make the proof readable; they do not replace UiPath orchestration.",
    },
    {
        "kind": "evidence",
        "eyebrow": "05 · SCENARIO 2A",
        "title": "Same green business systems. Missing signal means verify, not close.",
        "image": ART / "evidence_packet_E002_desktop.png",
        "caption": "The agent recommends closure. Policy preserves that recommendation, links it to the final decision, and routes the case to telemetry verification.",
        "points": [
            "Business state stays green.",
            "Closure is blocked by missing authoritative service evidence.",
            "The SLA clock stays alive until telemetry is fresh.",
        ],
    },
    {
        "kind": "evidence",
        "eyebrow": "06 · SCENARIO 2B",
        "title": "Same green business systems. Contradiction means human review.",
        "image": ART / "evidence_packet_E004_desktop.png",
        "caption": "The business fixture is intentionally the same. Fresh authoritative telemetry says the service is not live, so the route escalates.",
        "points": [
            "Only the authoritative evidence condition changes.",
            "Policy rejects closure and opens human review.",
            "The reviewer sees the contradiction and closure guard.",
        ],
    },
    {
        "kind": "comparison",
        "eyebrow": "07 · ONE VARIABLE, DIFFERENT ROUTE",
        "title": "Missing evidence and contradicting evidence are different operational problems.",
        "left": ("Missing or stale", ["Business systems are green", "Authoritative telemetry is absent or too old", "Allowed route: verify telemetry", "Operational posture: controlled retry"]),
        "right": ("Contradicting", ["Business systems are green", "Fresh authoritative telemetry says not live", "Allowed route: human review", "Operational posture: exception handling"]),
        "callout": "This is the central proof: same business-green fixture, different authoritative signal, different route.",
    },
    {
        "kind": "evidence",
        "eyebrow": "08 · AGENT USEFULNESS BOUNDARY",
        "title": "The LLM contributes structured interpretation, not final authority.",
        "image": ART / "evidence_packet_E003_adversarial_desktop.png",
        "caption": "An adversarial interpretation run turns closure pressure and unresolved risk into structured disagreement for policy.",
        "points": [
            "Advocate and skeptic outputs validate against the same schema.",
            "Disagreement is policy input, not prose parsing.",
            "Policy still owns closure eligibility.",
        ],
    },
    {
        "kind": "evals",
        "eyebrow": "09 · EVAL AND LEARNING LOOP",
        "title": "The eval suite keeps the agent useful without weakening closure policy.",
        "rows": [
            ("E-001", "Aligned closure accepted"),
            ("E-002/E-003", "Missing or stale evidence blocks closure"),
            ("E-004", "Source contradiction escalates"),
            ("E-005/E-006", "Notes and pressure influence routing without overriding truth"),
            ("E-007/E-009", "Invalid output and override persistence are tested"),
            ("E-008", "Policy improvement remains pending human approval"),
        ],
        "callout": "Validated path: local tests plus Test Manager manual representation, not claimed automated Test Cloud execution.",
    },
    {
        "kind": "feedback",
        "eyebrow": "10 · PRODUCT FEEDBACK AWARD",
        "title": "The strongest product ask is a Maestro Case Human-Review Readiness Check.",
        "lead": "The pieces exist across Maestro, Actions, packages, Data Fabric, and Test Manager. Builders need one preflight and auditability contract that proves the human-review path will work before runtime.",
        "cards": [
            ("Before runtime", "Check services, roles, Actions, package/feed version, task bindings, and evidence fields."),
            ("During runtime", "Expose reviewer field rendering, task lifecycle, and structured return status."),
            ("After runtime", "Reconstruct agent event, policy decision, human action, versions, and audit payload in one view."),
        ],
        "callout": "This is framed as actionable product feedback, not a complaint list.",
    },
    {
        "kind": "coding",
        "eyebrow": "11 · CODING AGENT CONTRIBUTION",
        "title": "Codex helped build and verify the submission pack.",
        "rows": [
            ("Built", "Python core, schemas, policy engine, evidence packets, audit bundles."),
            ("Validated", "Unit tests, eval harness, UiPath CLI runbooks, submission sanity check."),
            ("Documented", "Platform proof map, product feedback, demo storyboard, coding-agent manifest."),
            ("Boundary", "Codex is build-time assistance, not runtime case authority."),
        ],
        "callout": "Runtime authority remains: deterministic policy, UiPath orchestration, and human review.",
    },
    {
        "kind": "close",
        "eyebrow": "12 · IMPACT",
        "title": "ClearPath Recovery makes agentic service recovery safer without making it toothless.",
        "metrics": [
            ("0", "closures without fresh authoritative service evidence"),
            ("2", "exception paths proven visibly"),
            ("9/9", "eval scenarios represented in Test Manager"),
            ("67", "local tests passing in submission check"),
        ],
        "callout": "Agents interpret. Policy decides. Maestro routes. Humans own exceptions.",
    },
]


def add_textbox(slide, x, y, w, h, text="", size=18, color=NEAR_BLACK, bold=False, font=SERIF, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return box


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PARCHMENT


def footer(slide, index):
    add_textbox(slide, 0.55, 7.05, 3.2, 0.18, "SERVICE RECOVERY", 8, STONE, font=MONO)
    add_textbox(slide, 12.05, 7.05, 0.55, 0.18, f"{index:02d}", 8, STONE, font=MONO, align=PP_ALIGN.RIGHT)


def title_block(slide, eyebrow, title, lead=None):
    add_textbox(slide, 0.55, 0.38, 11.8, 0.25, eyebrow, 8.5, STONE, font=MONO)
    add_textbox(slide, 0.55, 0.70, 11.8, 0.85, title, 27, NEAR_BLACK)
    if lead:
        add_textbox(slide, 0.58, 1.54, 10.8, 0.58, lead, 13.2, OLIVE)


def card(slide, x, y, w, h, title, body, accent=BRAND):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = IVORY
    shape.line.color.rgb = BORDER
    add_textbox(slide, x + 0.18, y + 0.14, w - 0.36, 0.28, title, 14, accent)
    add_textbox(slide, x + 0.18, y + 0.54, w - 0.36, h - 0.62, body, 11.2, DARK_WARM)


def callout(slide, text):
    add_textbox(slide, 0.62, 6.36, 11.6, 0.48, text, 13.2, BRAND)


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, spec in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        set_bg(slide)

        if spec["kind"] == "cover":
            add_textbox(slide, 1.2, 2.15, 10.9, 0.78, spec["title"], 43, NEAR_BLACK, align=PP_ALIGN.CENTER)
            add_textbox(slide, 2.0, 3.05, 9.3, 0.55, spec["subtitle"], 17, OLIVE, align=PP_ALIGN.CENTER)
            add_textbox(slide, 2.4, 4.18, 8.5, 0.28, spec["meta"], 9.5, STONE, font=MONO, align=PP_ALIGN.CENTER)
            slide.shapes.add_shape(1, Inches(4.15), Inches(3.82), Inches(5.0), Inches(0.03)).fill.solid()
            slide.shapes[-1].fill.fore_color.rgb = BRAND
            continue

        title_block(slide, spec["eyebrow"], spec["title"], spec.get("lead"))
        kind = spec["kind"]

        if kind in {"problem", "casefit"}:
            for idx, (t, b) in enumerate(spec["cards"]):
                card(slide, 0.65 + idx * 4.05, 2.32, 3.6, 2.15, t, b)
            callout(slide, spec["callout"])
        elif kind == "architecture":
            x0 = 0.62
            for idx, (num, t, b) in enumerate(spec["steps"]):
                x = x0 + idx * 2.48
                add_textbox(slide, x, 2.08, 0.75, 0.35, num, 20, BRAND)
                add_textbox(slide, x, 2.58, 2.05, 0.34, t, 14.5, NEAR_BLACK)
                add_textbox(slide, x, 3.05, 2.0, 1.15, b, 10.8, DARK_WARM)
                if idx < 4:
                    add_textbox(slide, x + 2.02, 2.68, 0.3, 0.25, "→", 18, STONE)
            callout(slide, spec["callout"])
        elif kind == "map":
            y = 1.86
            for name, body in spec["rows"]:
                add_textbox(slide, 0.72, y, 2.3, 0.3, name, 13.2, BRAND)
                add_textbox(slide, 3.05, y, 8.95, 0.34, body, 11.2, DARK_WARM)
                y += 0.67
            callout(slide, spec["callout"])
        elif kind == "evidence":
            slide.shapes.add_picture(str(spec["image"]), Inches(0.62), Inches(1.82), width=Inches(7.4))
            add_textbox(slide, 8.48, 1.90, 4.0, 0.75, spec["caption"], 12, OLIVE)
            y = 3.05
            for point in spec["points"]:
                add_textbox(slide, 8.58, y, 3.82, 0.35, "• " + point, 12.4, DARK_WARM)
                y += 0.55
        elif kind == "comparison":
            left_title, left_items = spec["left"]
            right_title, right_items = spec["right"]
            card(slide, 0.78, 2.0, 5.65, 3.8, left_title, "\n".join(left_items), GREEN)
            card(slide, 6.90, 2.0, 5.65, 3.8, right_title, "\n".join(right_items), RED)
            callout(slide, spec["callout"])
        elif kind == "evals":
            y = 1.82
            for sid, body in spec["rows"]:
                add_textbox(slide, 0.72, y, 1.0, 0.28, sid, 12.4, BRAND, font=MONO)
                add_textbox(slide, 1.85, y, 9.8, 0.28, body, 12.3, DARK_WARM)
                y += 0.52
            callout(slide, spec["callout"])
        elif kind == "feedback":
            for idx, (t, b) in enumerate(spec["cards"]):
                card(slide, 0.78 + idx * 4.1, 2.25, 3.65, 2.3, t, b)
            callout(slide, spec["callout"])
        elif kind == "coding":
            y = 2.0
            for name, body in spec["rows"]:
                add_textbox(slide, 0.82, y, 1.35, 0.34, name, 13.5, BRAND)
                add_textbox(slide, 2.2, y, 9.45, 0.34, body, 12.6, DARK_WARM)
                y += 0.68
            callout(slide, spec["callout"])
        elif kind == "close":
            for idx, (num, label) in enumerate(spec["metrics"]):
                x = 0.88 + idx * 3.05
                add_textbox(slide, x, 2.25, 2.25, 0.55, num, 31, BRAND, align=PP_ALIGN.CENTER)
                add_textbox(slide, x, 3.0, 2.25, 0.72, label, 11.4, DARK_WARM, align=PP_ALIGN.CENTER)
            callout(slide, spec["callout"])

        footer(slide, i)

    prs.core_properties.title = "ClearPath Recovery"
    prs.core_properties.subject = "UiPath AgentHack Maestro Case presentation deck"
    prs.core_properties.author = "Arshdeep Singh"
    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(OUT / "governed_service_recovery_uipath_agenthack.pptx")


def make_official_template_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide():
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        return slide

    slide = new_slide()
    add_textbox(slide, 1.35, 2.15, 10.6, 0.72, "ClearPath Recovery", 40, NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_textbox(
        slide,
        2.2,
        3.0,
        8.9,
        0.68,
        "A UiPath Maestro Case workflow that prevents telecom service cases from closing until authoritative evidence proves recovery.",
        16,
        OLIVE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(slide, 2.5, 4.08, 8.4, 0.26, "Track 1 · Maestro Case · Governed Telecom Recovery", 9.5, STONE, font=MONO, align=PP_ALIGN.CENTER)

    slide = new_slide()
    title_block(slide, "02 · PROBLEM AND SOLUTION", "Problem statement and proposed solution")
    card(
        slide,
        0.78,
        1.95,
        5.65,
        3.55,
        "Problem",
        "Telecom service recovery can look complete in CRM, order, billing, and support systems while the customer still lacks working service. The danger is wrongful closure, repeat contact, SLA exposure, and an audit trail that cannot explain why the case closed.",
        BRAND,
    )
    card(
        slide,
        6.9,
        1.95,
        5.65,
        3.55,
        "Solution",
        "ClearPath Recovery uses agents to interpret ambiguous evidence into structured signals, deterministic policy to decide whether closure is allowed, Maestro Case to route dynamic exceptions, Action Center for human review, and Data Fabric/Orchestrator for durable audit proof.",
        BRAND,
    )
    callout(slide, "Agents interpret. Policy decides. Maestro routes. Humans own exceptions.")
    footer(slide, 2)

    slide = new_slide()
    title_block(slide, "03 · BENEFITS AND TECHNOLOGIES", "Benefits and technologies used")
    rows = [
        ("End user", "Telecom service recovery teams, activation/restoration owners, exception reviewers, audit/compliance teams."),
        ("User department", "Customer operations, network operations, field operations, service assurance, compliance."),
        ("Industries", "Telecom, broadband, managed connectivity, enterprise service operations."),
        ("UiPath products used", "Maestro Case, Action Center, Orchestrator, Data Fabric, Test Manager, UiPath CLI."),
        ("Other integrations / APIs", "Simulated CRM/order, billing, network telemetry, inventory, dispatch, support notes; optional Gemini/Vertex LLM interpreter."),
    ]
    y = 1.76
    for label, body in rows:
        add_textbox(slide, 0.78, y, 2.15, 0.32, label, 12.4, BRAND)
        add_textbox(slide, 3.02, y, 9.35, 0.36, body, 12.2, DARK_WARM)
        y += 0.58
    add_textbox(
        slide,
        0.78,
        4.92,
        11.3,
        0.65,
        "Benefits: fewer wrongful closures, lower repeat contacts, better SLA control, safer agent adoption, and an auditable chain from raw agent recommendation to policy decision to human action.",
        13.2,
        BRAND,
    )
    footer(slide, 3)

    slide = new_slide()
    title_block(slide, "04 · SOLUTION ARCHITECTURE", "Solution architecture")
    steps = [
        ("Evidence", "CRM, billing, telemetry, inventory, dispatch, notes"),
        ("Agent", "Structured interpretation of ambiguous text"),
        ("Policy", "Closure guard and allowed route"),
        ("Maestro Case", "Dynamic stage routing and lifecycle"),
        ("Human review", "High-impact exception ownership"),
    ]
    for idx, (name, body) in enumerate(steps):
        x = 0.72 + idx * 2.45
        add_textbox(slide, x, 2.05, 0.72, 0.4, f"{idx+1:02d}", 19, BRAND)
        add_textbox(slide, x, 2.65, 2.05, 0.28, name, 14, NEAR_BLACK)
        add_textbox(slide, x, 3.06, 1.95, 0.9, body, 10.8, DARK_WARM)
        if idx < len(steps) - 1:
            add_textbox(slide, x + 1.92, 2.78, 0.3, 0.25, "→", 17, STONE)
    callout(slide, "Track 1 fit: unpredictable evidence-driven routes, human-in-the-loop decisions, exception handling, visibility, and auditability.")
    footer(slide, 4)

    slide = new_slide()
    title_block(slide, "05 · DEMO PROOF", "Running proof and repository")
    slide.shapes.add_picture(str(ART / "evidence_packet_E004_desktop.png"), Inches(0.78), Inches(1.65), width=Inches(6.7))
    add_textbox(
        slide,
        7.9,
        1.82,
        4.5,
        0.75,
        "Core proof: same green business systems, different authoritative evidence, different route.",
        15,
        BRAND,
    )
    add_textbox(
        slide,
        8.02,
        3.0,
        4.1,
        1.35,
        "Missing/stale evidence routes to telemetry verification. Fresh contradiction routes to human review. The raw agent recommendation, final policy decision, reviewer packet, and audit chain are visible.",
        12.5,
        DARK_WARM,
    )
    add_textbox(slide, 8.02, 5.05, 4.0, 0.35, "GitHub: github.com/Arshgill01/uipath-agenthack-service-recovery", 10.5, STONE, font=MONO)
    footer(slide, 5)

    prs.core_properties.title = "ClearPath Recovery - Official Template Version"
    prs.core_properties.subject = "UiPath AgentHack submission deck"
    prs.core_properties.author = "Arshdeep Singh"
    prs.save(OUT / "clearpath_recovery_official_template_version.pptx")


def html_escape(value: str) -> str:
    return value.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def make_html():
    css = """
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>Governed Service Recovery</title>
<style>
@page { size: 280mm 158mm; margin: 0; background: #f5f4ed; }
* { box-sizing: border-box; }
:root { --parchment:#f5f4ed; --ivory:#faf9f5; --brand:#1B365D; --near:#141413; --warm:#3d3d3a; --olive:#504e49; --stone:#6b6a64; --border:#e8e6dc; --red:#b82217; --green:#1d6b57; }
body { margin:0; color:var(--near); font-family: Charter, Georgia, Palatino, serif; background:var(--parchment); }
.slide { width:280mm; height:158mm; padding:16mm 20mm; break-after:page; position:relative; background:var(--parchment); overflow:hidden; }
.cover { display:flex; flex-direction:column; align-items:center; justify-content:center; text-align:center; }
.cover h1 { font-size:38pt; line-height:1.06; font-weight:500; margin:0 0 12pt; letter-spacing:-0.5pt; }
.cover .sub { font-size:14pt; color:var(--olive); max-width:54ch; line-height:1.45; }
.meta,.eyebrow,.foot,.num { font-family: Menlo, Consolas, monospace; color:var(--stone); letter-spacing:1.5pt; text-transform:uppercase; font-size:9pt; }
.eyebrow { margin-bottom:2pt; }
h2 { font-size:24pt; line-height:1.16; font-weight:500; margin:0 0 11pt; max-width:920px; }
.lead { color:var(--olive); font-size:12pt; line-height:1.5; max-width:860px; margin-bottom:14pt; }
.cards { display:grid; grid-template-columns:repeat(3,1fr); gap:14pt; margin-top:18pt; }
.card { background:var(--ivory); border:0.5pt solid var(--border); border-radius:7pt; padding:13pt 15pt; min-height:37mm; }
.card h3 { color:var(--brand); font-size:15pt; margin:0 0 7pt; font-weight:500; }
.card p,.row p,.point { color:var(--warm); font-size:11pt; line-height:1.48; margin:0; }
.co { position:absolute; left:20mm; right:20mm; bottom:12mm; color:var(--brand); font-size:11.5pt; line-height:1.5; font-weight:500; }
.foot { position:absolute; left:20mm; bottom:10mm; }
.num { position:absolute; right:20mm; bottom:10mm; }
.steps { display:grid; grid-template-columns:repeat(5,1fr); gap:12pt; margin-top:23pt; }
.step .n { color:var(--brand); font-size:24pt; }
.step h3 { font-size:15pt; margin:7pt 0; font-weight:500; }
.rows { margin-top:12pt; }
.row { display:grid; grid-template-columns:44mm 1fr; border-bottom:0.3pt solid var(--border); padding:7pt 0; }
.row b { color:var(--brand); font-size:12pt; font-weight:500; }
.evidence { display:grid; grid-template-columns: 1.65fr 1fr; gap:18pt; margin-top:10pt; align-items:start; }
.evidence img { width:100%; border:0.5pt solid var(--border); border-radius:6pt; }
.caption { color:var(--olive); font-size:11pt; line-height:1.45; margin-bottom:12pt; }
.compare { display:grid; grid-template-columns:1fr 1fr; gap:22pt; margin-top:18pt; }
.compare .left h3 { color:var(--green); } .compare .right h3 { color:var(--red); }
.metric-grid { display:grid; grid-template-columns:repeat(4,1fr); gap:15pt; margin-top:32pt; text-align:center; }
.metric .big { color:var(--brand); font-size:34pt; line-height:1; }
.metric p { color:var(--warm); font-size:11pt; line-height:1.4; }
ul { margin:0; padding-left:16pt; } li { color:var(--warm); font-size:11pt; margin:0 0 8pt; line-height:1.4; }
</style></head><body>
"""
    body = [css]
    for i, spec in enumerate(SLIDES, 1):
        if spec["kind"] == "cover":
            body.append(f"""<section class="slide cover"><h1>{html_escape(spec['title'])}</h1><div class="sub">{html_escape(spec['subtitle'])}</div><div style="height:18pt"></div><div class="meta">{html_escape(spec['meta'])}</div></section>""")
            continue
        body.append(f"""<section class="slide"><div class="eyebrow">{html_escape(spec['eyebrow'])}</div><h2>{html_escape(spec['title'])}</h2>""")
        if spec.get("lead"):
            body.append(f"""<p class="lead">{html_escape(spec['lead'])}</p>""")
        kind = spec["kind"]
        if kind in {"problem", "casefit", "feedback"}:
            body.append('<div class="cards">')
            for t, b in spec["cards"]:
                body.append(f"""<div class="card"><h3>{html_escape(t)}</h3><p>{html_escape(b)}</p></div>""")
            body.append("</div>")
        elif kind == "architecture":
            body.append('<div class="steps">')
            for n, t, b in spec["steps"]:
                body.append(f"""<div class="step"><div class="n">{n}</div><h3>{html_escape(t)}</h3><p class="point">{html_escape(b)}</p></div>""")
            body.append("</div>")
        elif kind in {"map", "evals", "coding"}:
            body.append('<div class="rows">')
            for a, b in spec["rows"]:
                body.append(f"""<div class="row"><b>{html_escape(a)}</b><p>{html_escape(b)}</p></div>""")
            body.append("</div>")
        elif kind == "evidence":
            rel = spec["image"].relative_to(OUT).as_posix() if False else "../../demo/artifacts/" + spec["image"].name
            body.append(f"""<div class="evidence"><img src="{rel}"><div><p class="caption">{html_escape(spec['caption'])}</p><ul>""")
            for p in spec["points"]:
                body.append(f"<li>{html_escape(p)}</li>")
            body.append("</ul></div></div>")
        elif kind == "comparison":
            body.append('<div class="compare">')
            for cls, side in (("left", spec["left"]), ("right", spec["right"])):
                title, items = side
                body.append(f"""<div class="card {cls}"><h3>{html_escape(title)}</h3><ul>""")
                for item in items:
                    body.append(f"<li>{html_escape(item)}</li>")
                body.append("</ul></div>")
            body.append("</div>")
        elif kind == "close":
            body.append('<div class="metric-grid">')
            for n, label in spec["metrics"]:
                body.append(f"""<div class="metric"><div class="big">{html_escape(n)}</div><p>{html_escape(label)}</p></div>""")
            body.append("</div>")
        if spec.get("callout"):
            body.append(f"""<div class="co">{html_escape(spec['callout'])}</div>""")
        body.append(f"""<div class="foot">SERVICE RECOVERY</div><div class="num">{i:02d}</div></section>""")
    body.append("</body></html>")
    OUT.mkdir(parents=True, exist_ok=True)
    (OUT / "governed_service_recovery_kami_deck.html").write_text("\n".join(body), encoding="utf-8")


def make_guide():
    lines = [
        "# Organizer Template Adaptation Guide",
        "",
        "Official template link supplied by Devpost: https://bit.ly/3R0MsHU",
        "",
        "Use `governed_service_recovery_uipath_agenthack.pptx` as the finished editable deck. If the official Google Slides template must be used exactly, paste this content into the organizer slides while preserving the organizer title/body placeholders.",
        "",
        "A compact 5-slide companion deck matching the organizer template structure is generated as `clearpath_recovery_official_template_version.pptx`.",
        "",
        "## Slide Mapping",
        "",
    ]
    for i, spec in enumerate(SLIDES, 1):
        lines.append(f"### {i}. {spec['title']}")
        lines.append(f"- Layout: {spec['kind']}")
        if spec.get("eyebrow"):
            lines.append(f"- Eyebrow: {spec['eyebrow']}")
        if spec.get("subtitle"):
            lines.append(f"- Subtitle: {spec['subtitle']}")
        if spec.get("lead"):
            lines.append(f"- Lead: {spec['lead']}")
        if spec.get("callout"):
            lines.append(f"- Bottom callout: {spec['callout']}")
        if spec.get("image"):
            lines.append(f"- Image: `{spec['image'].relative_to(ROOT)}`")
        lines.append("")
    lines.extend(
        [
            "## Video Use",
            "",
            "Use slides 1, 3, 4, and 12 as quick framing in the video. Spend most of the video on the live/proof screens: E-002, E-004, E-003 adversarial, Test Manager, product feedback, Codex proof, and `scripts/run_submission_check.sh` output.",
            "",
            "## Claim Boundaries",
            "",
            "- Do not claim real telecom integration.",
            "- Do not claim automated Test Cloud execution.",
            "- Do not claim generated Action Center UI is the judge-facing proof surface.",
            "- Do not claim LLMs or Codex can close cases or override policy.",
        ]
    )
    (OUT / "organizer_template_adaptation_guide.md").write_text("\n".join(lines), encoding="utf-8")


def main():
    make_deck()
    make_official_template_deck()
    make_html()
    make_guide()
    print(f"Wrote deck artifacts to {OUT}")


if __name__ == "__main__":
    main()
